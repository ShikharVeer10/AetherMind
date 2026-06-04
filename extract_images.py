import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def main():
    pptx_path = r"C:\Users\shikh\Downloads\E-Commerce-Recommendation-System-23BCE9278.pptx"
    output_dir = Path("output/extracted_images")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation(pptx_path)
    print(f"Loaded presentation: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}")
    
    extracted_count = 0
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        img_idx = 1
        
        # Helper to extract recursively from groups
        def extract_from_shape(shape):
            nonlocal img_idx, extracted_count
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in getattr(shape, "shapes", []):
                    extract_from_shape(child)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    ext = image.ext if hasattr(image, "ext") else "png"
                    filename = f"slide_{slide_num}_image_{img_idx}.{ext}"
                    filepath = output_dir / filename
                    filepath.write_bytes(blob)
                    print(f"Extracted: {filename} ({len(blob)} bytes)")
                    img_idx += 1
                    extracted_count += 1
                except Exception as e:
                    print(f"Failed to extract image on slide {slide_num}: {e}")

        for shape in slide.shapes:
            extract_from_shape(shape)
            
    print(f"Extraction completed. Total images extracted: {extracted_count}")

if __name__ == "__main__":
    main()
