from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from models.document_model import (
    DocumentModel,
    DocumentElementModel,
    SlideModel,
    PositionModel,
    StyleModel,
    ParagraphModel,
    RunModel,
)


class PPTExtractor:
    def __init__(self, pptx_file_path: str):
        self.pptx_file_path = Path(pptx_file_path)
        self.presentation = Presentation(pptx_file_path)

    def extract_document(self) -> DocumentModel:
        extracted_slides = []
        for slide_index, slide in enumerate(self.presentation.slides):
            extracted_slide = self.extract_slide(
                slide=slide, slide_number=slide_index + 1
            )
            extracted_slides.append(extracted_slide)

        # Extract presentation-level metadata
        pres_meta = self._extract_presentation_metadata()

        return DocumentModel(
            document_name=Path(self.pptx_file_path).name,
            document_type="ppt",
            total_slides=len(extracted_slides),
            slides=extracted_slides,
            presentation_metadata=pres_meta,
        )

    def _extract_presentation_metadata(self) -> dict:
        """Extract top-level presentation metadata (author, dimensions, etc.)"""
        meta: dict = {}
        try:
            core = self.presentation.core_properties
            if core.author:
                meta["author"] = core.author
            if core.title:
                meta["presentation_title"] = core.title
            if core.subject:
                meta["subject"] = core.subject
            if core.created:
                meta["created"] = str(core.created)
            if core.modified:
                meta["modified"] = str(core.modified)
        except Exception:
            pass
        try:
            meta["slide_width"] = float(self.presentation.slide_width)
            meta["slide_height"] = float(self.presentation.slide_height)
        except Exception:
            pass
        return meta
    def extract_slide(self, slide, slide_number: int) -> SlideModel:
        extracted_elements = []
        slide_title: Optional[str] = None

        # 1. Try to get the real title from placeholders
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text_frame.text.strip() or None
        except Exception:
            pass

        for index, shape in enumerate(slide.shapes):
            elements = self._extract_shape_recursive(
                shape=shape,
                slide_number=slide_number,
                shape_index=index,
                z_order=index,
                prefix=f"slide_{slide_number}",
            )
            for element in elements:
                extracted_elements.append(element)

        # 2. Fallback: if no placeholder title, use the first element text
        if slide_title is None:
            for element in extracted_elements:
                if element.text:
                    slide_title = element.text.strip()
                    break

        # Extract slide background color
        slide_bg_color = None
        try:
            background = slide.background
            if background and background.fill:
                if background.fill.type == 1 and background.fill.fore_color: # solid fill
                    slide_bg_color = self._get_safe_color_hex(background.fill.fore_color)
        except Exception:
            pass

        return SlideModel(
            slide_number=slide_number,
            title=slide_title,
            elements=extracted_elements,
            background_color=slide_bg_color,
        )

    def _extract_shape_recursive(
        self,
        shape,
        slide_number: int,
        shape_index: int,
        prefix: str,
        z_order: int = 0,
    ) -> List[DocumentElementModel]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_elements = []
            group_shapes = getattr(shape, "shapes", [])
            for child_idx, child_shape in enumerate(group_shapes):
                child_prefix = f"{prefix}_group_{shape_index}"
                child_elements.extend(
                    self._extract_shape_recursive(
                        shape=child_shape,
                        slide_number=slide_number,
                        shape_index=child_idx,
                        z_order=z_order + child_idx,
                        prefix=child_prefix,
                    )
                )
            return child_elements
        element = self._extract_single_shape(
            shape=shape,
            element_id=f"{prefix}_shape_{shape_index}",
            z_order=z_order,
        )
        return [element] if element is not None else []

    def _extract_single_shape(
        self, shape, element_id: str, z_order: int = 0
    ) -> Optional[DocumentElementModel]:
        paragraphs: List[ParagraphModel] = []
        full_text: Optional[str] = None

        if shape.has_text_frame:
            paragraph_texts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue

                runs = []
                for run in paragraph.runs:
                    run_text = run.text
                    if not run_text:
                        continue
                    run_model = RunModel(
                        text=run_text,
                        bold=bool(run.font.bold),
                        italic=bool(run.font.italic),
                        font_size=(
                            float(run.font.size.pt) if run.font.size else None
                        ),
                        font_name=run.font.name,
                        font_color=self._get_safe_color_hex(run.font.color),
                    )
                    runs.append(run_model)

                # Extract paragraph alignment
                alignment_name = None
                if paragraph.alignment is not None:
                    try:
                        alignment_name = paragraph.alignment.name
                    except Exception:
                        alignment_name = str(paragraph.alignment)

                paragraphs.append(
                    ParagraphModel(
                        level=paragraph.level,
                        text=para_text,
                        runs=runs,
                        alignment=alignment_name,
                    )
                )

                paragraph_texts.append(para_text)

            if paragraph_texts:
                full_text = "\n".join(paragraph_texts)

        x, y, w, h = self._safe_position(shape)
        position = PositionModel(
            x=x,
            y=y,
            width=w,
            height=h,
        )

        style = self._extract_text_style(shape)
        element_type = self._get_shape_type(shape)
        metadata = self._extract_shape_metadata(shape, element_type)
        metadata["z_order"] = z_order
        table_md: Optional[str] = None
        if element_type == "table":
            table_md = self.extract_table_as_markdown(shape)
        return DocumentElementModel(
            element_id=element_id,
            element_type=element_type,
            text=full_text,
            paragraphs=paragraphs,
            position=position,
            style=style,
            shape_type=str(shape.shape_type),
            metadata=metadata,
            table_markdown=table_md,
        )

    def _get_shape_type(self, shape) -> str:
        st = shape.shape_type

        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "text_box"

        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            name = shape.name.lower()
            # Detect arrows from shape name
            if "arrow" in name:
                return "arrow"
            # Detect process/decision shapes common in flowcharts
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            if auto_shape_type is not None:
                auto_name = str(auto_shape_type).lower()
                if 'arrow' in auto_name:
                    return "arrow"
                if 'diamond' in auto_name or 'decision' in auto_name:
                    return "shape"  # decision diamond
                if 'chevron' in auto_name or 'pentagon' in auto_name:
                    return "shape"  # process chevron
            return "shape"

        if st == MSO_SHAPE_TYPE.GROUP:
            return "group"

        if st == MSO_SHAPE_TYPE.TABLE:
            return "table"

        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"

        if st == MSO_SHAPE_TYPE.CHART:
            return "chart"

        if st == MSO_SHAPE_TYPE.FREEFORM:
            return "freeform"

        if st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"

        connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
        if connector_type is not None and st == connector_type:
            return "connector"
        if st == MSO_SHAPE_TYPE.LINE:
            return "connector"

        return "unknown"


    def _extract_shape_metadata(self, shape, element_type: str) -> dict:
        metadata: dict = {}
        metadata["name"] = getattr(shape, "name", "")
        try:
            metadata["rotation"] = float(shape.rotation or 0)
        except Exception:
            metadata["rotation"] = 0
        try:
            metadata["visible"] = bool(getattr(shape, "visible", True))
        except Exception:
            metadata["visible"] = True
        try:
            metadata["is_placeholder"] = bool(getattr(shape, "is_placeholder", False))
        except Exception:
            metadata["is_placeholder"] = False

        if element_type == "table":
            metadata["table_data"] = self._extract_table_data(shape)
        if element_type == "chart":
            metadata["chart_data"] = self._extract_pptx_chart_data(shape)
        if element_type == "image":
            try:
                metadata["__image_bytes"] = shape.image.blob
            except Exception:
                try:
                    # Fallback to direct relationship blob access for custom content type images (e.g. image/jpg)
                    blip = shape._element.blipFill.blip
                    embed_key = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                    link_key = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link'
                    rId = blip.get(embed_key) or blip.get(link_key)
                    rel = shape.part.rels[rId]
                    target_part = rel.target_part
                    if hasattr(target_part, "blob"):
                        metadata["__image_bytes"] = target_part.blob
                    elif hasattr(target_part, "_blob"):
                        metadata["__image_bytes"] = target_part._blob
                    else:
                        metadata["__image_bytes"] = None
                except Exception:
                    metadata["__image_bytes"] = None

        if shape.has_text_frame:
            metadata["bullet_hierarchy"] = self._extract_bullet_hierarchy(shape)
        if element_type == "connector":
            metadata["connector_endpoints"] = (
                self._extract_connector_endpoints(shape)
            )

        # High-fidelity rendering additions: Auto Shape geometry type and borders
        try:
            if hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
                metadata["auto_shape_type"] = shape.auto_shape_type.name
        except Exception:
            pass


        try:
            if hasattr(shape, "line") and shape.line:
                border_color = self._get_safe_color_hex(shape.line.color)
                if border_color:
                    metadata["border_color"] = border_color
                if shape.line.width:
                    metadata["border_width"] = float(shape.line.width.pt)
        except Exception:
            pass


        return metadata


    def _extract_table_data(self, shape) -> list:
        rows_out = []
        if not shape.has_table:
            return rows_out
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_out.append(cells)
        return rows_out

    def _extract_pptx_chart_data(self, shape) -> Optional[dict]:
        if not hasattr(shape, "chart") or shape.chart is None:
            return None
        try:
            chart = shape.chart
            title = ""
            if chart.has_title:
                title = chart.chart_title.text_frame.text

            series_data = []
            for s in chart.series:
                series_data.append({
                    "name": s.name,
                    "values": list(s.values) if hasattr(s, "values") else []
                })

            categories = []
            if len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, "categories"):
                    categories = [str(cat) for cat in plot.categories]

            chart_type_name = "unknown"
            if hasattr(chart, "chart_type"):
                chart_type_name = str(chart.chart_type)

            return {
                "title": title,
                "chart_type": chart_type_name,
                "series": series_data,
                "categories": categories
            }
        except Exception:
            return None

    def _extract_bullet_hierarchy(self, shape) -> list:
        items = []
        if not shape.has_text_frame:
            return items
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                items.append({"level": paragraph.level, "text": text})
        return items

    def _extract_connector_endpoints(self, shape) -> dict:
        endpoints: dict = {
            "begin_x": None,
            "begin_y": None,
            "end_x": None,
            "end_y": None,
        }
        # Prefer python-pptx Connector properties (reliable)
        if hasattr(shape, "begin_x"):
            try:
                endpoints["begin_x"] = float(shape.begin_x)
                endpoints["begin_y"] = float(shape.begin_y)
                endpoints["end_x"] = float(shape.end_x)
                endpoints["end_y"] = float(shape.end_y)
                return endpoints
            except Exception:
                pass
        # Fallback to XML parsing
        try:
            sp_element = shape._element
            xfrm = sp_element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
            )
            if xfrm is not None:
                off = xfrm.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
                )
                ext = xfrm.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}ext"
                )
                if off is not None:
                    endpoints["begin_x"] = float(off.get("x", 0))
                    endpoints["begin_y"] = float(off.get("y", 0))
                if off is not None and ext is not None:
                    endpoints["end_x"] = (
                        float(off.get("x", 0)) + float(ext.get("cx", 0))
                    )
                    endpoints["end_y"] = (
                        float(off.get("y", 0)) + float(ext.get("cy", 0))
                    )
        except Exception:
            pass
        return endpoints

    def _safe_position(self, shape) -> tuple[float, float, float, float]:
        try:
            return (
                float(shape.left),
                float(shape.top),
                float(shape.width),
                float(shape.height),
            )
        except Exception:
            pass

        try:
            sp_element = shape._element
            xfrm = sp_element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
            )
            if xfrm is None:
                return 0.0, 0.0, 0.0, 0.0
            off = xfrm.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
            )
            ext = xfrm.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}ext"
            )
            x = float(off.get("x", 0)) if off is not None else 0.0
            y = float(off.get("y", 0)) if off is not None else 0.0
            w = float(ext.get("cx", 0)) if ext is not None else 0.0
            h = float(ext.get("cy", 0)) if ext is not None else 0.0
            return x, y, w, h
        except Exception:
            return 0.0, 0.0, 0.0, 0.0


    def extract_table_as_markdown(self, shape) -> str:
        if not shape.has_table:
            return ""

        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
        body_lines = []
        for row in rows[1:]:
            body_lines.append("| " + " | ".join(row) + " |")

        parts = [header, separator] + body_lines
        return "\n".join(parts)

    def _get_safe_color_hex(self, color_obj) -> Optional[str]:
        """Safely extract hex color, avoiding ValueError on scheme colors."""
        if not color_obj:
            return None
        try:
            if color_obj.type == MSO_COLOR_TYPE.RGB and color_obj.rgb:
                return self.convert_rgb_to_hex(color_obj.rgb)
        except Exception:
            pass
        return None

    def _extract_text_style(self, shape) -> StyleModel:
        font_size: Optional[float] = None
        font_name: Optional[str] = None
        is_bold = False
        is_italic = False
        text_color = None
        background_color = None

        try:
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                if paragraphs and paragraphs[0].runs:
                    first_run = paragraphs[0].runs[0]
                    font = first_run.font
                    if font.size:
                        font_size = float(font.size.pt)
                    if font.name:
                        font_name = font.name
                    if font.bold:
                        is_bold = True
                    if font.italic:
                        is_italic = True
                    text_color = self._get_safe_color_hex(font.color)
        except Exception:
            pass

        try:
            if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                background_color = self.convert_rgb_to_hex(shape.fill.fore_color.rgb)
        except Exception:
            pass

        return StyleModel(
            font_size=font_size,
            font_name=font_name,
            bold=is_bold,
            italic=is_italic,
            text_color=text_color,
            background_color=background_color,
        )

    def sort_elements_by_reading_order(self, slide):
        """Sort elements top-to-bottom, then left-to-right."""
        return sorted(
            slide.elements,
            key=lambda e: (e.position.y, e.position.x),
        )

    @staticmethod
    def convert_rgb_to_hex(rgb_color: RGBColor) -> str:
        return "#{:02x}{:02x}{:02x}".format(
            rgb_color[0], rgb_color[1], rgb_color[2]
        )
