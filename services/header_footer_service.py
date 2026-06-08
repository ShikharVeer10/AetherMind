from typing import Optional
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from models.document_model import HeaderFooterModel


# Top 12% of a standard 16:9 slide (6,858,000 EMU height)
_HEADER_Y_CUTOFF = 820_000
# Bottom 12% — anything below this Y is considered footer region
_FOOTER_Y_CUTOFF = 6_030_000


class HeaderFooterService:
    """
    Reads header/footer placeholders from three layers:
        1. The slide itself
        2. The slide's layout
        3. The slide master

    Additionally scans regular text shapes positioned in the extreme top
    or bottom of the slide as secondary header/footer sources.
    """

    def extract(self, slide) -> HeaderFooterModel:
        header_text: Optional[str] = None
        footer_text: Optional[str] = None
        slide_number_text: Optional[str] = None
        date_text: Optional[str] = None

        if slide is None or not hasattr(slide, "placeholders"):
            return HeaderFooterModel(
                header_text=header_text,
                footer_text=footer_text,
                slide_number_text=slide_number_text,
                date_text=date_text,
            )

        # --- Layer 1: Slide placeholders ---
        for ph in slide.placeholders:
            header_text, footer_text, slide_number_text, date_text = (
                self._check_placeholder(
                    ph, header_text, footer_text,
                    slide_number_text, date_text,
                )
            )

        # --- Layer 2: Slide layout placeholders ---
        try:
            layout = slide.slide_layout
            for ph in layout.placeholders:
                header_text, footer_text, slide_number_text, date_text = (
                    self._check_placeholder(
                        ph, header_text, footer_text,
                        slide_number_text, date_text,
                    )
                )
        except Exception:
            pass

        # --- Layer 3: Slide master placeholders ---
        try:
            master = slide.slide_layout.slide_master
            for ph in master.placeholders:
                header_text, footer_text, slide_number_text, date_text = (
                    self._check_placeholder(
                        ph, header_text, footer_text,
                        slide_number_text, date_text,
                    )
                )
        except Exception:
            pass

        # --- Layer 4: Positional fallback for header/footer from text shapes ---
        # --- Layer 4: Positional fallback for header/footer from text shapes ---
        header_text, footer_text = self._positional_header_footer(slide,header_text,footer_text)

        confidentiality_label = self._detect_confidentiality(footer_text)

        return HeaderFooterModel(
            header_text=header_text,
            footer_text=footer_text,
            slide_number_text=slide_number_text,
            date_text=date_text,
            confidentiality_label=confidentiality_label,
        )

    @staticmethod
    def _positional_header_footer(
    slide,
    existing_header: Optional[str],
    existing_footer: Optional[str],
    existing_slide_number: Optional[str] = None,
    ):

        header = existing_header
        footer = existing_footer
        slide_number = existing_slide_number

        for shape in slide.shapes:

            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()

            if not text:
                continue

            import re

            if slide_number is None:
                if (bottom > _FOOTER_Y_CUTOFF and re.fullmatch(r"\d+", text)):
                    slide_number = text

            try:
                top = float(shape.top)
                bottom = top + float(shape.height)

            except Exception:
                continue

            if header is None and top < _HEADER_Y_CUTOFF:
                header = text

            elif footer is None and bottom > _FOOTER_Y_CUTOFF:
                footer = text
        return(header,footer,slide_number)


    @staticmethod
    def _check_placeholder(
        ph,
        header: Optional[str],
        footer: Optional[str],
        slide_num: Optional[str],
        date: Optional[str],
    ):
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            return header, footer, slide_num, date

        text = ""
        if ph.has_text_frame:
            text = ph.text_frame.text.strip()

        if ph_type == PP_PLACEHOLDER_TYPE.FOOTER and footer is None:
            footer = text or None

        elif ph_type == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER and slide_num is None:
            slide_num = text or None

        elif ph_type == PP_PLACEHOLDER_TYPE.DATE and date is None:
            date = text or None

        elif (
            ph_type == PP_PLACEHOLDER_TYPE.TITLE
            and header is None
            and ph.top is not None
            and float(ph.top) < 500_000
        ):
            header = text or None

        return header, footer, slide_num, date
    
    @staticmethod
    def _detect_confidentiality(
        text: str | None
    ) -> str | None:

        if not text:
            return None

        lower = text.lower()

        keywords = [
            "confidential",
            "internal use only",
            "proprietary",
            "draft",
            "restricted"
        ]

        for keyword in keywords:
            if keyword in lower:
                return keyword

        return None
